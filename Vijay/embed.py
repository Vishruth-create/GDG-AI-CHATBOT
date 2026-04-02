import groq
import base64
import os
import fitz
import io
from PIL import Image
from pptx import Presentation
from docx import Document
import openpyxl 
from langchain_text_splitters import RecursiveCharacterTextSplitter
from qdrant_client.models import Distance, VectorParams, PointStruct
from utilities import clean_text, is_garbage_text, setup_qdrant, load_model, verify_insert
from config import collection_name
from config import collection_name, config

vector_size     = config.vector_size     
batch_size      = config.batch_size      
chunk_size      = config.chunk_size         
chunk_overlap   = config.chunk_overlap 
from dotenv import load_dotenv
load_dotenv()
groq_client = groq.Groq(api_key=os.getenv("groq_api_key"))

# Keep your original function unchanged
def analyze_image_with_groq(image: Image.Image) -> str:
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    img_base64 = base64.b64encode(buffer.getvalue()).decode("utf-8")

    try:
        response = groq_client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[{
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_base64}"}},
                    {"type": "text", "text": "Extract all text from this image. If no text, describe what you see in two or more sentences"}
                ]
            }]
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"groq error: {e}")
        return ""


def concatenate_images_vertically(images: list[Image.Image]) -> Image.Image:
    """Stack images vertically into single image"""
    if not images:
        return None
    
    max_width = max(img.width for img in images)
    resized = []
    
    for img in images:
        if img.width != max_width:
            ratio = max_width / img.width
            new_height = int(img.height * ratio)
            img = img.resize((max_width, new_height), Image.LANCZOS)
        resized.append(img)
    
    total_height = sum(img.height for img in resized)
    combined = Image.new('RGB', (max_width, total_height), 'white')
    
    y_offset = 0
    for img in resized:
        combined.paste(img, (0, y_offset))
        y_offset += img.height
    
    return combined


def load_ppt(file_path):
    filename = os.path.basename(file_path)
    prs = Presentation(file_path)
    pages = []
    
    # Collect all slides data first
    all_slides_data = []
    for slide_num, slide in enumerate(prs.slides):
        slide_text = []
        slide_images = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        slide_text.append(line)
            elif shape.shape_type == 13:
                try:
                    img_bytes = shape.image.blob
                    slide_images.append(Image.open(io.BytesIO(img_bytes)))
                except Exception as e:
                    print(f"Skipped image: {e}")
        
        all_slides_data.append({
            'slide_num': slide_num,
            'text': slide_text,
            'images': slide_images
        })
    
    # Process in batches of 4 slides
    for batch_start in range(0, len(all_slides_data), 4):
        batch = all_slides_data[batch_start:batch_start + 4]
        
        # Collect all images from these 4 slides
        batch_images = []
        for slide_data in batch:
            batch_images.extend(slide_data['images'])
        
        # Concatenate and analyze once for 4 slides
        batch_image_text = ""
        if batch_images:
            combined_img = concatenate_images_vertically(batch_images)
            batch_image_text = analyze_image_with_groq(combined_img)
        
        # Create pages for each slide in batch
        for slide_data in batch:
            full_text = " ".join(slide_data['text']) + " " + batch_image_text
            full_text = " ".join(full_text.split())
            
            if len(full_text.strip()) < 3:
                continue
            
            cleaned = clean_text(full_text)
            if len(cleaned.strip()) >= 3:
                pages.append({
                    "text": cleaned,
                    "page_num": slide_data['slide_num'] + 1,
                    "source": filename
                })
    
    return pages

def extract_images_text_from_pdf_page(doc, page):
    image_texts = []
    image_list = page.get_images(full=True)
    
    for img in image_list:
        xref = img[0]#unique id for imag
        try:
            base_image  = doc.extract_image(xref) #extract raw image byte from pdf
            image_bytes = base_image["image"]
            image    = Image.open(io.BytesIO(image_bytes))
            ocr_text = analyze_image_with_groq(image)
            
            if ocr_text.strip():
                image_texts.append(ocr_text.strip())
                
        except Exception as e:
            print(f"Skipping pdf image error: {e}")
            continue
            
    return " ".join(image_texts)

def load_pdf(file_path):
    filename = os.path.basename(file_path)
    doc      = fitz.open(file_path)
    pages    = []
    skipped  = 0

    for page_num in range(len(doc)):
        page = doc[page_num] 
        raw_text = page.get_text()
        image_text = extract_images_text_from_pdf_page(doc, page)
        full_text = raw_text + " " + image_text
        if is_garbage_text(full_text):
            skipped += 1
            continue
        cleaned = clean_text(full_text)
        if len(cleaned) < 50:
            skipped += 1
            continue

        pages.append({
            "text"    : cleaned,
            "page_num": page_num + 1,
            "source"  : filename
        })
        
    doc.close()
    return pages

def load_image(file_path):
    filename = os.path.basename(file_path)
    image    = Image.open(file_path)
    text     = analyze_image_with_groq(image)
    cleaned  = clean_text(text)
    if len(cleaned) < 5:
        return []
    return [{"text": cleaned, "page_num": 1, "source": filename}]

def load_docx(file_path):
    filename = os.path.basename(file_path)
    doc      = Document(file_path)
    pages    = []
    para_batch = []
    page_num   = 1

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            para_batch.append(text)

        # every 20 paragraphs treat as one page
        if len(para_batch) >= 20:
            full_text = " ".join(para_batch)
            cleaned   = clean_text(full_text)

            if not is_garbage_text(full_text) and len(cleaned) >= 50:
                pages.append({
                    "text"    : cleaned,
                    "page_num": page_num,
                    "source"  : filename
                })
                page_num += 1

            para_batch = []

    #remaining paragraphs
    if para_batch:
        full_text = " ".join(para_batch)
        cleaned   = clean_text(full_text)

        if not is_garbage_text(full_text) and len(cleaned) >= 50:
            pages.append({
                "text"    : cleaned,
                "page_num": page_num,
                "source"  : filename
            })
    return pages



#load xlsx file row by row using openpyxl
def load_xlsx(file_path):
    filename = os.path.basename(file_path)
    wb       = openpyxl.load_workbook(file_path, data_only=True)
    pages    = []

    print(f"opening excel file '{filename}' {len(wb.sheetnames)} sheets")

    for sheet_num, sheet_name in enumerate(wb.sheetnames):
        ws         = wb[sheet_name]
        sheet_rows = []

        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(
                str(cell) for cell in row
                if cell is not None and str(cell).strip()
            )
            if row_text.strip():
                sheet_rows.append(row_text)

        full_text = " ".join(sheet_rows)

        if len(full_text.strip()) < 10:
            continue
        cleaned = clean_text(full_text)

        pages.append({
            "text"    : cleaned,
            "page_num": sheet_num + 1,
            "source"  : filename
        })
    return pages

# detect file type and call the right loader
def load_file(file_path):
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return load_pdf(file_path)
    elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"]:
        return load_image(file_path)
    elif ext in [".pptx", ".ppt"]:
        return load_ppt(file_path)
    elif ext in [".xlsx", ".xls"]:
        return load_xlsx(file_path)
    elif ext in [".docx", ".doc"]: 
        return load_docx(file_path)
    else:
        print(f"Unsupported file type: {ext}")
        return []

def make_chunks(pages):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size     =chunk_size,
        chunk_overlap  =chunk_overlap,
        separators     =["\n\n", "\n", ". ", " ", ""],
        length_function=len
    )

    all_chunks = []
    chunk_id   = 0

    for page in pages:
        chunks = splitter.split_text(page["text"])

        for chunk in chunks:
            chunk = chunk.strip()
            if len(chunk) < 80:
                continue

            all_chunks.append({
                "chunk_text": chunk,
                "page_num"  : page["page_num"],
                "source"    : page["source"],
                "chunk_id"  : chunk_id
            })
            chunk_id += 1

    print(f"Total chunks: {len(all_chunks)}")
    return all_chunks

def embed_chunks(chunks, model):
    texts = [chunk["chunk_text"] for chunk in chunks]
    print(f"Embedding {len(texts)} chunks")

    vectors = model.encode(
        texts,
        batch_size          =batch_size,
        show_progress_bar   =True,
        convert_to_numpy    =True,
        normalize_embeddings=True
    )

    return vectors

def create_collection(client):
    existing = [c.name for c in client.get_collections().collections]

    if collection_name in existing:
        client.delete_collection(collection_name=collection_name)
        print(f"Old collection '{collection_name}' deleted")

    client.create_collection(
        collection_name=collection_name,
        vectors_config =VectorParams(
            size    =vector_size,
            distance=Distance.COSINE
        )
    )
    print(f"Collection '{collection_name}' created")
    return True

def insert_to_qdrant(chunks, vectors, client):
    points = []

    for i, (chunk, vector) in enumerate(zip(chunks, vectors)):
        points.append(PointStruct(
            id     = i,
            vector = vector.tolist(),
            payload = {
                "chunk_text": chunk["chunk_text"],
                "page_num"  : chunk["page_num"],
                "source"    : chunk["source"],
                "chunk_id"  : chunk["chunk_id"]
            }
        ))

    client.upsert(
        collection_name=collection_name,
        points=points
    )
    print(f"Inserted {len(points)} points")


def main_pipeline(file_paths):
    client = setup_qdrant()
    create_collection(client)

    # accept single file or list of files
    if isinstance(file_paths, str):
        file_paths = [file_paths]

    all_pages  = []

    for file_path in file_paths:
        print(f"processing: {file_path}")
        pages = load_file(file_path)
        all_pages.extend(pages)


    chunks  = make_chunks(all_pages)
    model   = load_model()
    vectors = embed_chunks(chunks, model)
    insert_to_qdrant(chunks, vectors, client)
    verify_insert(client)

    return client, model

if __name__ == "__main__":
    files = [
        # "GDG Inductions 2026.pdf",
        #add files to process
        
    ]
    client, model = main_pipeline(files)